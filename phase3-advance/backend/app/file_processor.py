"""File Processing Module - Professional美化增强 - Word/Excel/PPT generation."""
import os, uuid, shutil, re
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Any, Optional, Tuple
from loguru import logger

from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font as XlFont, PatternFill, Alignment as XlAlignment, Border, Side, numbers
from openpyxl.chart import BarChart, LineChart, PieChart, Reference
from openpyxl.utils import get_column_letter

from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn, nsdecls
from docx.oxml import parse_xml

from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import pandas as pd
from PIL import Image

try: import PyPDF2
except ImportError: PyPDF2 = None
try: import pdfplumber
except ImportError: pdfplumber = None

from app.config import FILES_DIR, UPLOAD_DIR, GENERATED_DIR, BACKUP_DIR


class WordTheme:
    THEMES = {
        "academic": {"title":"1A3A5C","h1":"1A3A5C","h2":"2C5F8A","accent":"C49A2B","cover_bg":"1A3A5C","cover_text":"FFFFFF"},
        "business": {"title":"1E3A5F","h1":"1E3A5F","h2":"34495E","accent":"2980B9","cover_bg":"1E3A5F","cover_text":"FFFFFF"},
        "modern":  {"title":"2563EB","h1":"1E40AF","h2":"3B82F6","accent":"F59E0B","cover_bg":"2563EB","cover_text":"FFFFFF"},
        "warm_friendly": {"title":"E8943A","h1":"D4772A","h2":"F0A050","accent":"4A90D9","cover_bg":"E8943A","cover_text":"FFFFFF"},
    }
    @classmethod
    def get(cls, name="academic"): return cls.THEMES.get(name, cls.THEMES["academic"])

class XlTheme:
    THEMES = {
        "professional": {"hdr_bg":"2F5496","hdr_fg":"FFFFFF","alt1":"FFFFFF","alt2":"D6E4F0","border":"B0B0B0","chart_accent":["2F5496","2E75B6","C55A11","548235","BF8F00"]},
        "dark":         {"hdr_bg":"333333","hdr_fg":"FFFFFF","alt1":"FFFFFF","alt2":"F2F2F2","border":"CCCCCC","chart_accent":["333333","555555","888888","AAAAAA","CCCCCC"]},
        "corporate":    {"hdr_bg":"002060","hdr_fg":"FFFFFF","alt1":"FFFFFF","alt2":"E6EDF5","border":"B8CCE4","chart_accent":["002060","2E75B6","C55A11","548235","BF8F00"]},
    }
    @classmethod
    def get(cls, name="professional"): return cls.THEMES.get(name, cls.THEMES["professional"])


class FileProcessor:
    SUPPORTED_TYPES = {"pdf":"PDF","docx":"Word","xlsx":"Excel","csv":"CSV","txt":"Text","png":"Image","jpg":"Image","jpeg":"Image","bmp":"Image","gif":"Image"}
    MAX_FILE_SIZE = 20*1024*1024

    def __init__(self):
        for d in [UPLOAD_DIR, GENERATED_DIR, BACKUP_DIR]: d.mkdir(parents=True, exist_ok=True)
    def validate_file(self, fn, sz):
        ext = fn.rsplit(".",1)[-1].lower() if "." in fn else ""
        if ext not in self.SUPPORTED_TYPES: return False, f"不支持格式：.{ext}"
        if sz > self.MAX_FILE_SIZE: return False, "文件过大"
        return True, ""
    def save_upload(self, content, fn):
        ext = fn.rsplit(".",1)[-1].lower() if "." in fn else "bin"
        p = UPLOAD_DIR / f"{uuid.uuid4().hex}.{ext}"; p.write_bytes(content); return str(p)
    def backup_file(self, fp):
        src = Path(fp)
        if not src.exists(): return ""
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        tgt = BACKUP_DIR / f"{src.stem}_backup_{ts}{src.suffix}"
        shutil.copy2(src, tgt); return str(tgt)
    def _hex(self, h):
        h = h.lstrip("#")
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    # ═══════════════════════════════════════════
    # EXCEL - Enhanced with #2F5496 header, horizontal borders, number format
    # ═══════════════════════════════════════════

    def generate_excel(self, sheet_name, headers, rows, title="", chart_type="", theme_name="professional"):
        theme = XlTheme.get(theme_name)
        wb = Workbook(); ws = wb.active; ws.title = sheet_name or "Sheet1"
        hdr_font = XlFont(name="微软雅黑", bold=True, size=12, color=theme["hdr_fg"])
        hdr_fill = PatternFill(start_color=theme["hdr_bg"], end_color=theme["hdr_bg"], fill_type="solid")
        hdr_align = XlAlignment(horizontal="center", vertical="center", wrap_text=True)
        data_font = XlFont(name="微软雅黑", size=11, color="333333")
        data_align_left = XlAlignment(horizontal="left", vertical="center")
        data_align_center = XlAlignment(horizontal="center", vertical="center")
        data_align_right = XlAlignment(horizontal="right", vertical="center")
        h_border = Border(bottom=Side(style="hair", color=theme["border"]))
        alt1 = PatternFill(start_color=theme["alt1"], end_color=theme["alt1"], fill_type="solid")
        alt2 = PatternFill(start_color=theme["alt2"], end_color=theme["alt2"], fill_type="solid")
        title_font = XlFont(name="微软雅黑", bold=True, size=16, color=theme["hdr_bg"])

        hs = 1
        if title:
            ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=len(headers))
            ws.cell(row=1, column=1, value=title).font = title_font
            ws.cell(row=1, column=1).alignment = XlAlignment(horizontal="center", vertical="center")
            ws.row_dimensions[1].height = 35; hs = 2

        for ci, h in enumerate(headers, 1):
            c = ws.cell(row=hs, column=ci, value=h)
            c.font = hdr_font; c.fill = hdr_fill; c.alignment = hdr_align; c.border = h_border
        ws.row_dimensions[hs].height = 30

        for ri, rd in enumerate(rows, hs+1):
            for ci, v in enumerate(rd, 1):
                c = ws.cell(row=ri, column=ci, value=v); c.font = data_font; c.border = h_border
                c.fill = alt1 if (ri-hs)%2==1 else alt2
                if isinstance(v, (int, float)): c.alignment = data_align_right
                elif self._is_date_str(str(v)): c.alignment = data_align_center
                else: c.alignment = data_align_left
            ws.row_dimensions[ri].height = 22

        for ci in range(1, len(headers)+1):
            mx = len(str(headers[ci-1]))*2+4
            for ri in range(hs+1, hs+1+len(rows)):
                val = str(ws.cell(row=ri, column=ci).value or "")
                w = sum(2 if '\u4e00'<=c<='\u9fff' else 1 for c in val)
                mx = max(mx, w+4)
            ws.column_dimensions[get_column_letter(ci)].width = min(max(mx,8),30)

        ws.freeze_panes = f"A{hs+1}"
        if rows: ws.auto_filter.ref = f"A{hs}:{get_column_letter(len(headers))}{hs+len(rows)}"

        if chart_type and rows:
            self._add_xl_chart(wb, ws, chart_type, hs, len(rows), len(headers), theme)

        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe = (title or sheet_name or "table").replace("/","_")[:30]
        fp = GENERATED_DIR / f"{safe}_{ts}.xlsx"; wb.save(str(fp))
        return str(fp)

    def _is_date_str(self, s): return bool(re.match(r'^\d{4}[-/年]\d{1,2}[-/月]\d{1,2}[日]?$', s))

    def _add_xl_chart(self, wb, ws, ct, sr, nr, nc, theme):
        try:
            chart_cls = {"bar":BarChart,"line":LineChart,"pie":PieChart}
            if ct not in chart_cls: return
            ch = chart_cls[ct](); ch.title="数据分析图表"; ch.style=10
            if ct!="pie": ch.y_axis.title="数值"; ch.x_axis.title="类别"
            dr = Reference(ws, min_col=1, min_row=sr, max_col=nc, max_row=sr+nr)
            cr = Reference(ws, min_col=1, min_row=sr+1, max_row=sr+nr)
            ch.add_data(dr, titles_from_data=True); ch.set_categories(cr)
            ws.add_chart(ch, f"{get_column_letter(nc+2)}{sr}")
        except: pass

    # ═══════════════════════════════════════════
    # WORD - Enhanced cover with color block + decorative line
    # ═══════════════════════════════════════════

    def generate_word(self, title, content, author="小石头", theme_name="academic", doc_type="report"):
        theme = WordTheme.get(theme_name)
        doc = Document()
        s = doc.sections[0]
        s.page_width=Cm(21.0); s.page_height=Cm(29.7)
        s.top_margin=Cm(2.54); s.bottom_margin=Cm(2.54)
        s.left_margin=Cm(3.18); s.right_margin=Cm(3.18)
        st = doc.styles["Normal"]
        st.font.name="宋体"; st.font.size=Pt(12); st.font.color.rgb=RGBColor(0x33,0x33,0x33)
        st.paragraph_format.line_spacing=1.5; st.paragraph_format.first_line_indent=Cm(0.74); st.paragraph_format.space_after=Pt(6)
        rpr = st.element.rPr
        if rpr is None: rpr=parse_xml(f'<w:rPr {nsdecls("w")}></w:rPr>'); st.element.append(rpr)
        rf = rpr.find(qn("w:rFonts"))
        if rf is None: rf=parse_xml(f'<w:rFonts {nsdecls("w")}/>'); rpr.append(rf)
        rf.set(qn("w:eastAsia"),"宋体")

        self._add_word_cover(doc, title, author, theme, doc_type)
        doc.add_page_break()
        toc=doc.add_paragraph(); toc.alignment=WD_ALIGN_PARAGRAPH.CENTER
        r=toc.add_run("目  录"); r.font.size=Pt(18); r.font.bold=True
        r.font.color.rgb=self._hex(theme["title"]); r.font.name="微软雅黑"
        doc.add_paragraph(); doc.add_page_break()
        self._add_word_content(doc, content, theme)

        # Header & Footer
        for sec in doc.sections:
            h=sec.header; h.is_linked_to_previous=False
            hp=h.paragraphs[0] if h.paragraphs else h.add_paragraph()
            hp.alignment=WD_ALIGN_PARAGRAPH.CENTER
            rh=hp.add_run(title); rh.font.size=Pt(9); rh.font.color.rgb=RGBColor(0x99,0x99,0x99); rh.font.name="宋体"
            # Footer with page number
            f=sec.footer; f.is_linked_to_previous=False
            fp=f.paragraphs[0] if f.paragraphs else f.add_paragraph()
            fp.alignment=WD_ALIGN_PARAGRAPH.CENTER
            r1=fp.add_run("第"); r1.font.size=Pt(9); r1.font.color.rgb=RGBColor(0x99,0x99,0x99)
            r2=fp.add_run(); r2._r.append(parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="begin"/>'))
            r3=fp.add_run(); r3._r.append(parse_xml(f'<w:instrText {nsdecls("w")} xml:space="preserve"> PAGE </w:instrText>'))
            r4=fp.add_run(); r4._r.append(parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="end"/>'))
            r5=fp.add_run("页"); r5.font.size=Pt(9); r5.font.color.rgb=RGBColor(0x99,0x99,0x99)

        ts=datetime.now().strftime("%Y%m%d_%H%M%S")
        safe=re.sub(r'[\\/*?:"<>|]','_',title)[:30]
        fp=GENERATED_DIR/f"{safe}_{ts}.docx"; doc.save(str(fp)); return str(fp)

    def _add_word_cover(self, doc, title, author, theme, doc_type):
        # Color bar at top
        bar=doc.add_paragraph(); bar.paragraph_format.first_line_indent=Cm(0); bar.paragraph_format.space_after=Pt(0)
        bar.paragraph_format.space_before=Pt(0)
        shd=parse_xml(f'<w:shd {nsdecls("w")} w:val="clear" w:color="auto" w:fill="{theme.get("cover_bg",theme["title"])}"/>')
        bar._p.get_or_add_pPr().append(shd)
        for _ in range(4): bar.add_run("\n").font.size=Pt(4)

        # Title
        tp=doc.add_paragraph(); tp.alignment=WD_ALIGN_PARAGRAPH.CENTER
        tp.paragraph_format.space_before=Pt(80); tp.paragraph_format.first_line_indent=Cm(0)
        tr=tp.add_run(title); tr.font.size=Pt(30); tr.font.bold=True
        tr.font.color.rgb=self._hex(theme["title"]); tr.font.name="微软雅黑"

        # Decorative line
        sp=doc.add_paragraph(); sp.alignment=WD_ALIGN_PARAGRAPH.CENTER; sp.paragraph_format.first_line_indent=Cm(0)
        sr=sp.add_run("━"*50); sr.font.color.rgb=self._hex(theme["accent"]); sr.font.size=Pt(14)

        # Subtitle
        tn={"report":"研究报告","resume":"个人简历","letter":"正式信函","proposal":"项目提案"}
        sb=doc.add_paragraph(); sb.alignment=WD_ALIGN_PARAGRAPH.CENTER; sb.paragraph_format.first_line_indent=Cm(0)
        sb.paragraph_format.space_after=Pt(80)
        ss=sb.add_run(f"—— {tn.get(doc_type,'文档')} ——"); ss.font.size=Pt(18)
        ss.font.color.rgb=self._hex(theme["h2"]); ss.font.name="微软雅黑"

        # Author & Date
        inf=doc.add_paragraph(); inf.alignment=WD_ALIGN_PARAGRAPH.CENTER; inf.paragraph_format.first_line_indent=Cm(0)
        if author: inf.add_run(f"作者：{author}\n\n").font.size=Pt(13)
        inf.add_run(f"日期：{datetime.now().strftime('%Y年%m月%d日')}").font.size=Pt(13)

    def _add_word_content(self, doc, content, theme):
        lines=content.split("\n"); i=0
        while i<len(lines):
            line=lines[i].strip()
            if not line: i+=1; p=doc.add_paragraph(); p.paragraph_format.first_line_indent=Cm(0); continue
            if line.startswith("# ") and not line.startswith("## "):
                p=doc.add_paragraph(); p.paragraph_format.first_line_indent=Cm(0)
                p.paragraph_format.space_before=Pt(24); p.paragraph_format.space_after=Pt(12)
                r=p.add_run(line[2:]); r.font.size=Pt(22); r.font.bold=True; r.font.color.rgb=self._hex(theme["title"]); r.font.name="微软雅黑"
            elif line.startswith("## "):
                p=doc.add_paragraph(); p.paragraph_format.first_line_indent=Cm(0)
                p.paragraph_format.space_before=Pt(18); p.paragraph_format.space_after=Pt(8)
                r=p.add_run(line[3:]); r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=self._hex(theme["h1"]); r.font.name="微软雅黑"
            elif line.startswith("### "):
                p=doc.add_paragraph(); p.paragraph_format.first_line_indent=Cm(0)
                p.paragraph_format.space_before=Pt(12); p.paragraph_format.space_after=Pt(6)
                r=p.add_run(line[4:]); r.font.size=Pt(14); r.font.bold=True; r.font.color.rgb=self._hex(theme["h2"]); r.font.name="微软雅黑"
            elif line.startswith("- ") or line.startswith("* "):
                p=doc.add_paragraph(); p.style=doc.styles["List Bullet"]; p.paragraph_format.first_line_indent=Cm(0)
                r=p.add_run(line[2:]); r.font.name="宋体"; r.font.size=Pt(12)
            elif re.match(r'^\d+\.\s', line):
                p=doc.add_paragraph(); p.style=doc.styles["List Number"]; p.paragraph_format.first_line_indent=Cm(0)
                r=p.add_run(re.sub(r'^\d+\.\s','',line)); r.font.name="宋体"; r.font.size=Pt(12)
            elif line.startswith("|") and line.endswith("|"): i=self._add_word_table(doc, lines, i, theme)
            elif line.startswith("---"):
                p=doc.add_paragraph(); p.paragraph_format.first_line_indent=Cm(0)
                r=p.add_run("─"*60); r.font.color.rgb=RGBColor(0xcc,0xcc,0xcc); r.font.size=Pt(8)
            else: p=doc.add_paragraph(line)
            i+=1

    def _add_word_table(self, doc, lines, si, theme):
        rows=[]; i=si
        while i<len(lines) and lines[i].strip().startswith("|"):
            cells=[c.strip() for c in lines[i].strip().strip("|").split("|")]; rows.append(cells); i+=1
        dr=[r for r in rows if not all(re.match(r'^[-:]+$',c) for c in r)]
        if len(dr)<2: return i-1
        t=doc.add_table(rows=len(dr),cols=len(dr[0])); t.style="Table Grid"; t.alignment=WD_TABLE_ALIGNMENT.CENTER
        for ci,val in enumerate(dr[0]):
            c=t.cell(0,ci); c.text=val
            for p in c.paragraphs:
                p.alignment=WD_ALIGN_PARAGRAPH.CENTER; p.paragraph_format.first_line_indent=Cm(0)
                for rr in p.runs:
                    rr.font.size = Pt(10); rr.font.bold = True
                    rr.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); rr.font.name = '微软雅黑'
            c._tc.get_or_add_tcPr().append(parse_xml(f'<w:shd {nsdecls("w")} w:fill="{theme["title"]}"/>'))
        for ri in range(1,len(dr)):
            for ci,val in enumerate(dr[ri]):
                c=t.cell(ri,ci); c.text=val
                for p in c.paragraphs:
                    p.paragraph_format.first_line_indent=Cm(0)
                    for rr in p.runs:
                        rr.font.size = Pt(10); rr.font.name = '宋体'
                if ri%2==0: c._tc.get_or_add_tcPr().append(parse_xml(f'<w:shd {nsdecls("w")} w:fill="E8EDF2"/>'))
        doc.add_paragraph(); return i-1

    # ═══════════════════════════════════════════
    # PPT - Enhanced with card layout + decoration
    # ═══════════════════════════════════════════

    class PPTTheme:
        THEMES = {"academic_blue":{"bg":"1A3A5C","accent":"C49A2B","tlight":"FFFFFF","tdark":"1F2937"},
                  "elegant_gray":{"bg":"2D3436","accent":"0984E3","tlight":"FFFFFF","tdark":"2D3436"},
                  "modern_clean":{"bg":"FFFFFF","accent":"2563EB","tlight":"FFFFFF","tdark":"1E293B"}}
        @classmethod
        def get(cls,n="academic_blue"): return cls.THEMES.get(n, cls.THEMES["academic_blue"])

    def generate_ppt(self, title, slides_data, theme_name="academic_blue", author="", date_str=""):
        theme = self.PPTTheme.get(theme_name)
        prs = Presentation(); prs.slide_width=PptInches(13.333); prs.slide_height=PptInches(7.5)
        hrgb=lambda h: PptRGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))

        # ━━ COVER ━━
        sl=prs.slides.add_slide(prs.slide_layouts[6])
        sl.background.fill.solid(); sl.background.fill.fore_color.rgb=hrgb(theme["bg"])
        # Decorative circle (top-right)
        circ=sl.shapes.add_shape(9, PptInches(10.5), PptInches(-1.0), PptInches(3.0), PptInches(3.0))
        circ.fill.solid(); circ.fill.fore_color.rgb=hrgb(theme["accent"]); circ.line.fill.background()
        circ2=sl.shapes.add_shape(9, PptInches(0.5), PptInches(5.5), PptInches(1.5), PptInches(1.5))
        circ2.fill.solid(); circ2.fill.fore_color.rgb=hrgb(theme["accent"]); circ2.line.fill.background()
        # Title
        tb=sl.shapes.add_textbox(PptInches(1.5),PptInches(2.0),PptInches(10.3),PptInches(2.0))
        tf=tb.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.text=title
        p.font.size=PptPt(40); p.font.bold=True; p.font.color.rgb=hrgb(theme["tlight"]); p.font.name="微软雅黑"; p.alignment=PP_ALIGN.CENTER
        tb2=sl.shapes.add_textbox(PptInches(3.0),PptInches(4.3),PptInches(7.3),PptInches(0.5))
        p2=tb2.text_frame.paragraphs[0]; p2.text=f"—— {author or '答辩人'} · 毕业论文答辩"
        p2.font.size=PptPt(18); p2.font.color.rgb=hrgb(theme["accent"]); p2.font.name="微软雅黑"; p2.alignment=PP_ALIGN.CENTER

        # ━━ CONTENT SLIDES (card-style) ━━
        for sd in slides_data:
            ct=sd.get("content",[]); st=sd.get("title","")
            if isinstance(ct, str): ct=[ct]
            sl=prs.slides.add_slide(prs.slide_layouts[6])
            sl.background.fill.solid(); sl.background.fill.fore_color.rgb=hrgb(theme["tlight"])
            # Title bar with accent line
            bar=sl.shapes.add_shape(1,PptInches(0),PptInches(0),PptInches(13.333),PptInches(1.1))
            bar.fill.solid(); bar.fill.fore_color.rgb=hrgb(theme["bg"]); bar.line.fill.background()
            tt=sl.shapes.add_textbox(PptInches(0.8),PptInches(0.15),PptInches(11.5),PptInches(0.8))
            tp=tt.text_frame.paragraphs[0]; tp.text=st; tp.font.size=PptPt(28); tp.font.bold=True
            tp.font.color.rgb=hrgb(theme["tlight"]); tp.font.name="微软雅黑"
            # Card background (rounded rect)
            card=sl.shapes.add_shape(5, PptInches(0.8),PptInches(1.4),PptInches(11.7),PptInches(5.5))
            card.fill.solid(); card.fill.fore_color.rgb=hrgb(theme["tlight"]); card.line.color.rgb=hrgb(theme.get("accent","C49A2B")); card.line.width=Pt(1)
            # Content inside card
            cbox=sl.shapes.add_textbox(PptInches(1.3),PptInches(1.8),PptInches(10.7),PptInches(5.0))
            ctf=cbox.text_frame; ctf.word_wrap=True
            for j, item in enumerate(ct):
                p=ctf.paragraphs[0] if j==0 else ctf.add_paragraph()
                p.text=f"●  {item}" if not item.startswith("●") else item
                p.font.size=PptPt(20); p.font.color.rgb=hrgb(theme["tdark"]); p.font.name="微软雅黑"; p.space_after=PptPt(10)

        # ━━ CONCLUSION ━━
        sl=prs.slides.add_slide(prs.slide_layouts[6])
        sl.background.fill.solid(); sl.background.fill.fore_color.rgb=hrgb(theme["bg"])
        tb=sl.shapes.add_textbox(PptInches(1.5),PptInches(2.5),PptInches(10.3),PptInches(1.5))
        tf=tb.text_frame; tf.paragraphs[0].text="感谢聆听"; tf.paragraphs[0].font.size=PptPt(48)
        tf.paragraphs[0].font.bold=True; tf.paragraphs[0].font.color.rgb=hrgb(theme["tlight"]); tf.paragraphs[0].font.name="微软雅黑"; tf.paragraphs[0].alignment=PP_ALIGN.CENTER

        ts=datetime.now().strftime("%Y%m%d_%H%M%S"); safe=title.replace("/","_")[:30]
        fp=GENERATED_DIR/f"答辩_{safe}_{ts}.pptx"; prs.save(str(fp)); return str(fp)

    # ━━━━━━━ DATA CLEANING / PARSING ━━━━━━━
    def clean_data(self, fp, ops=None):
        ext=Path(fp).suffix.lower().lstrip(".")
        if ops is None: ops=["deduplicate","remove_nulls"]
        if ext in ("xlsx","csv"): return self._clean_tabular(fp, ext, ops)
        return {"cleaned_file":"","operations_performed":[]}
    def _clean_tabular(self, fp, ext, ops):
        try: df=pd.read_csv(fp,encoding="utf-8") if ext=="csv" else pd.read_excel(fp)
        except: df=pd.read_csv(fp,encoding="gbk")
        rb=len(df); res={"operations_performed":[],"rows_before":rb}
        if "deduplicate" in ops: df=df.drop_duplicates(); res["operations_performed"].append(f"去重{rb-len(df)}条")
        if "remove_nulls" in ops: df=df.dropna(how="all"); res["operations_performed"].append("删除空行")
        if "fill_missing" in ops: df=df.fillna("N/A"); res["operations_performed"].append("填充缺失")
        res["rows_after"]=len(df)
        ts=datetime.now().strftime("%Y%m%d_%H%M%S"); cp=GENERATED_DIR/f"cleaned_{ts}.{ext}"
        if ext=="csv": df.to_csv(cp,index=False,encoding="utf-8-sig")
        else: self._save_styled_xl(df,str(cp))
        res["cleaned_file"]=str(cp); return res
    def _save_styled_xl(self, df, fp):
        wb=Workbook(); ws=wb.active
        for ci,cn in enumerate(df.columns,1):
            ws.cell(row=1,column=ci,value=str(cn)).font=XlFont(name="微软雅黑",bold=True,size=12,color="FFFFFF")
            ws.cell(row=1,column=ci).fill=PatternFill(start_color="2F5496",end_color="2F5496",fill_type="solid")
        for ri,(_,row) in enumerate(df.iterrows(),2):
            for ci,v in enumerate(row,1):
                ws.cell(row=ri,column=ci,value=v).font=XlFont(name="微软雅黑",size=11)
                if ri%2==0: ws.cell(row=ri,column=ci).fill=PatternFill(start_color="D6E4F0",end_color="D6E4F0",fill_type="solid")
        wb.save(fp)
    def parse(self, fp):
        path=Path(fp); ext=path.suffix.lower().lstrip(".")
        result={"filename":path.name,"file_type":ext,"content":"","preview":"","rows":0}
        parsers={"pdf":self._p_pdf,"docx":self._p_docx,"xlsx":self._p_xlsx,"csv":self._p_csv,"txt":self._p_txt,"png":self._p_img,"jpg":self._p_img,"jpeg":self._p_img,"bmp":self._p_img,"gif":self._p_img}
        if ext in parsers:
            try: result.update(parsers[ext](fp))
            except Exception as e: result["content"]=f"[解析失败:{e}]"
        result["preview"]=result["content"][:500]; return result
    def _p_pdf(self,fp):
        t=[]
        if pdfplumber:
            try:
                with pdfplumber.open(fp) as pdf:
                    for p in pdf.pages:
                        x=p.extract_text()
                        if x: t.append(x)
            except: pass
        if not t and PyPDF2:
            try:
                with open(fp,"rb") as f:
                    for p in PyPDF2.PdfReader(f).pages:
                        x=p.extract_text()
                        if x: t.append(x)
            except: pass
        return {"content":"\n\n".join(t) if t else "[PDF解析库未安装]","rows":len(t)}
    def _p_docx(self,fp):
        doc=Document(fp); paras=[p.text for p in doc.paragraphs if p.text.strip()]
        return {"content":"\n".join(paras),"rows":len(paras)}
    def _p_xlsx(self,fp):
        try:
            dfs=pd.read_excel(fp,sheet_name=None); all_t=[]; total=0
            for n,df in dfs.items(): all_t.append(f"=== {n} ===\n{df.to_string(index=False)}"); total+=len(df)
            return {"content":"\n".join(all_t),"rows":total}
        except: return {"content":"[Excel解析失败]"}
    def _p_csv(self,fp):
        try: df=pd.read_csv(fp,encoding="utf-8")
        except: df=pd.read_csv(fp,encoding="gbk")
        return {"content":df.to_string(index=False),"rows":len(df)}
    def _p_txt(self,fp):
        for enc in ["utf-8","gbk"]:
            try: t=Path(fp).read_text(encoding=enc); return {"content":t,"rows":len(t.split("\n"))}
            except: continue
        return {"content":"[编码错误]","rows":0}
    def _p_img(self,fp):
        try: img=Image.open(fp); return {"content":f"[图片{img.size[0]}x{img.size[1]}]","rows":0}
        except: return {"content":"[图片解析失败]","rows":0}
    def modify_word(self,fp,mods):
        self.backup_file(fp); doc=Document(fp)
        if "append_text" in mods: doc.add_paragraph(mods["append_text"])
        if "replace_text" in mods:
            for old,new in mods["replace_text"].items():
                for p in doc.paragraphs:
                    if old in p.text:
                        for rr in p.runs:
                            if old in rr.text: rr.text=rr.text.replace(old,new)
        doc.save(fp); return fp
    def modify_excel(self,fp,mods):
        self.backup_file(fp); wb=load_workbook(fp); ws=wb.active
        if "append_rows" in mods:
            nr=ws.max_row+1
            for ri,row in enumerate(mods["append_rows"],nr):
                for ci,val in enumerate(row,1): ws.cell(row=ri,column=ci,value=val)
        if "update_cell" in mods:
            for ref,val in mods["update_cell"].items(): ws[ref]=val
        wb.save(fp); return fp